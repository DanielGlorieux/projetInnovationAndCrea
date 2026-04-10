"""
Script de création de la présentation PowerPoint sur l'algorithme IVF.
Basé sur les informations extraites du tutoriel TensorTeach.

La présentation peut être importée dans Canva via :
  https://www.canva.com/ -> "Importer" -> sélectionner le fichier .pptx

Prérequis :
    pip install python-pptx Pillow
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Couleurs du thème (inspirées du tutoriel TensorTeach - fond sombre, texte clair)
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)  # Bleu foncé profond
BG_SECTION = RGBColor(0x16, 0x21, 0x3E)  # Bleu marine
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)  # Bleu vif
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)  # Vert
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)  # Orange
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)  # Violet
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_YELLOW = RGBColor(0xFF, 0xD7, 0x00)

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output")
PPTX_PATH = os.path.join(OUTPUT_DIR, "presentation_ivf.pptx")

# Slide dimensions: 16:9 widescreen aspect ratio in inches
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Définit la couleur de fond d'une slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Ajoute un rectangle coloré à la slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Ajoute une zone de texte formatée à la slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, left, top, width, height, bullets, font_size=16,
                      color=TEXT_LIGHT, bullet_color=ACCENT_BLUE):
    """Ajoute une liste à puces formatée."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {bullet_text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.space_before = Pt(4)

    return txBox


def create_title_slide(prs):
    """Slide 1 : Titre de la présentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, BG_DARK)

    # Ligne décorative en haut
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.05), ACCENT_BLUE)

    # Titre principal
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Inverted File Index (IVF)", font_size=44, color=TEXT_WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Sous-titre
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
                 "Comprendre l'indexation dans les bases de données vectorielles",
                 font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    # Ligne de séparation
    add_shape_rect(slide, Inches(4.5), Inches(4.5), Inches(4), Inches(0.03), ACCENT_BLUE)

    # Source
    add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.8),
                 "Inspiré du tutoriel TensorTeach - Mastering Vector Databases",
                 font_size=16, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def create_intro_slide(prs):
    """Slide 2 : Introduction - Pourquoi IVF ?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Pourquoi l'IVF ?", font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_BLUE)

    # Problème
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.6),
                 "Le problème", font_size=22, color=ACCENT_ORANGE, bold=True)

    add_bullet_points(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(2.5), [
        "Les bases de données vectorielles stockent des millions de vecteurs",
        "La recherche exhaustive (brute force) est trop lente : O(n)",
        "Il faut un moyen d'accélérer la recherche",
        "Besoin d'un index efficace pour les vecteurs de haute dimension",
    ], font_size=15, color=TEXT_LIGHT)

    # Solution
    add_text_box(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.6),
                 "La solution : IVF", font_size=22, color=ACCENT_GREEN, bold=True)

    add_bullet_points(slide, Inches(6.8), Inches(2.4), Inches(5.5), Inches(2.5), [
        "Diviser l'espace vectoriel en clusters (groupes)",
        "Chercher uniquement dans les clusters les plus proches",
        "Réduire considérablement le nombre de comparaisons",
        "Compromis contrôlable entre vitesse et précision",
    ], font_size=15, color=TEXT_LIGHT)

    # Encadré récapitulatif
    box = add_shape_rect(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(1.5),
                         RGBColor(0x1E, 0x2D, 0x50), ACCENT_BLUE)
    add_text_box(slide, Inches(2), Inches(5.4), Inches(9), Inches(1.2),
                 "L'IVF (Inverted File Index) est un algorithme d'indexation qui partitionne "
                 "l'espace vectoriel en clusters via K-means, permettant de limiter la recherche "
                 "aux clusters les plus pertinents lors d'une requête.",
                 font_size=16, color=TEXT_YELLOW, alignment=PP_ALIGN.CENTER)


def create_indexing_overview_slide(prs):
    """Slide 3 : Vue d'ensemble de l'indexation IVF."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "IVF - Phase d'Indexation", font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.6),
                 "Comment l'IVF organise les données pour une recherche rapide",
                 font_size=18, color=TEXT_LIGHT)

    # Les 3 étapes
    steps = [
        ("1", "Clustering K-means",
         "Tous les vecteurs sont regroupés\nen 'nlist' clusters via l'algorithme\nK-means.",
         ACCENT_BLUE),
        ("2", "Centroïdes & Posting Lists",
         "Chaque cluster a un centroïde.\nUne posting list conserve les IDs\ndes vecteurs de ce cluster.",
         ACCENT_GREEN),
        ("3", "Ajout de nouveaux vecteurs",
         "Pour ajouter un vecteur, on trouve\nle centroïde le plus proche et on\nl'ajoute à sa posting list.",
         ACCENT_ORANGE),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.8 + i * 4.1)
        y = Inches(2.6)

        # Numéro circulaire
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Titre de l'étape
        add_text_box(slide, x, y + Inches(1), Inches(3.8), Inches(0.6),
                     title, font_size=20, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x, y + Inches(1.6), Inches(3.8), Inches(1.8),
                     desc, font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def create_clustering_slide(prs):
    """Slide 4 : Détail du clustering K-means."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Étape 1 : Clustering avec K-means", font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_BLUE)

    # Explication à gauche
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
                 "Principe", font_size=22, color=ACCENT_BLUE, bold=True)

    add_bullet_points(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(3), [
        "L'ensemble des vecteurs est partitionné en 'nlist' groupes (clusters)",
        "L'algorithme K-means est utilisé pour le clustering",
        "Chaque cluster est représenté par son centroïde (point central)",
        "Le paramètre 'nlist' contrôle le nombre de clusters",
        "Plus de clusters → des posting lists plus petites par cluster",
        "Moins de clusters → des posting lists plus grandes mais moins de clusters à vérifier",
    ], font_size=14, color=TEXT_LIGHT)

    # Illustration à droite (représentation schématique des clusters)
    cluster_box = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4),
                                 RGBColor(0x12, 0x1A, 0x30), ACCENT_BLUE)

    add_text_box(slide, Inches(7), Inches(2), Inches(5), Inches(0.5),
                 "Espace vectoriel 2D (simplifié)", font_size=14, color=TEXT_LIGHT,
                 alignment=PP_ALIGN.CENTER)

    # Dessiner des clusters schématiques
    cluster_colors = [
        (RGBColor(0xFF, 0x45, 0x45), "Cluster 1"),
        (RGBColor(0x45, 0xFF, 0x45), "Cluster 2"),
        (RGBColor(0x45, 0x45, 0xFF), "Cluster 3"),
        (RGBColor(0xFF, 0xFF, 0x45), "Cluster 4"),
    ]

    positions = [
        (Inches(7.3), Inches(2.8)),
        (Inches(9.8), Inches(2.8)),
        (Inches(7.3), Inches(4.2)),
        (Inches(9.8), Inches(4.2)),
    ]

    for (color, label), (cx, cy) in zip(cluster_colors, positions):
        # Cluster circle
        cluster = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(1.8), Inches(1.2))
        cluster.fill.solid()
        # Darken the cluster color to ~20% brightness for the fill background
        cluster.fill.fore_color.rgb = RGBColor(
            min(color[0] + 20, 255) // 5,
            min(color[1] + 20, 255) // 5,
            min(color[2] + 20, 255) // 5,
        )
        cluster.line.color.rgb = color
        cluster.line.width = Pt(2)

        # Centroid dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx + Inches(0.7), cy + Inches(0.4),
                                     Inches(0.3), Inches(0.3))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # Label
        add_text_box(slide, cx, cy + Inches(0.05), Inches(1.8), Inches(0.4),
                     label, font_size=11, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Note en bas
    box = add_shape_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1),
                         RGBColor(0x1E, 0x2D, 0x50), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), Inches(5.9), Inches(10.8), Inches(0.8),
                 "Le parametre 'nlist' determine le nombre de clusters. "
                 "Typiquement : nlist = sqrt(n) ou n est le nombre total de vecteurs.",
                 font_size=15, color=TEXT_YELLOW)


def create_posting_list_slide(prs):
    """Slide 5 : Centroïdes et Posting Lists."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Étape 2 : Centroïdes & Posting Lists",
                 font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_GREEN)

    add_bullet_points(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2), [
        "Chaque cluster possède un centroïde (point représentatif)",
        "Une 'posting list' associe chaque centroïde à la liste des IDs de ses vecteurs",
        "Structure similaire à un index inversé en recherche textuelle",
    ], font_size=15, color=TEXT_LIGHT)

    # Illustration posting list
    y_start = Inches(3.5)
    clusters_data = [
        ("Centroïde 1", ["ID: 5", "ID: 12", "ID: 23", "ID: 41"], ACCENT_BLUE),
        ("Centroïde 2", ["ID: 3", "ID: 8", "ID: 17"], ACCENT_GREEN),
        ("Centroïde 3", ["ID: 1", "ID: 9", "ID: 28", "ID: 35", "ID: 42"], ACCENT_ORANGE),
        ("Centroïde 4", ["ID: 7", "ID: 14", "ID: 31"], ACCENT_PURPLE),
    ]

    for i, (centroid_name, ids, color) in enumerate(clusters_data):
        x = Inches(0.8 + i * 3.1)

        # Centroid box
        cbox = add_shape_rect(slide, x, y_start, Inches(2.8), Inches(0.6), color)
        tf = cbox.text_frame
        tf.paragraphs[0].text = centroid_name
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Arrow
        add_text_box(slide, x, y_start + Inches(0.6), Inches(2.8), Inches(0.4),
                     "↓ Posting List", font_size=12, color=color,
                     alignment=PP_ALIGN.CENTER)

        # IDs
        for j, vid in enumerate(ids):
            id_box = add_shape_rect(slide, x + Inches(0.5), y_start + Inches(1.0 + j * 0.4),
                                    Inches(1.8), Inches(0.35),
                                    RGBColor(0x1E, 0x2D, 0x50), color)
            tf = id_box.text_frame
            tf.paragraphs[0].text = vid
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def create_adding_vector_slide(prs):
    """Slide 6 : Ajout d'un nouveau vecteur."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Étape 3 : Ajout d'un Nouveau Vecteur",
                 font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_ORANGE)

    # Process steps
    steps = [
        ("1. Calculer les distances",
         "Calculer la distance entre le nouveau\nvecteur et tous les centroïdes",
         ACCENT_BLUE),
        ("2. Trouver le plus proche",
         "Identifier le centroïde le plus\nproche du nouveau vecteur",
         ACCENT_GREEN),
        ("3. Ajouter à la posting list",
         "Ajouter l'ID du nouveau vecteur\nà la posting list du centroïde\nle plus proche",
         ACCENT_ORANGE),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.8 + i * 4.1)
        y = Inches(2)

        box = add_shape_rect(slide, x, y, Inches(3.8), Inches(2.8),
                             RGBColor(0x1E, 0x2D, 0x50), color)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.6),
                     title, font_size=18, color=color, bold=True)

        add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(3.4), Inches(1.2),
                     desc, font_size=14, color=TEXT_LIGHT)

        # Arrow between boxes
        if i < 2:
            add_text_box(slide, x + Inches(3.8), y + Inches(1), Inches(0.3), Inches(0.6),
                         "→", font_size=28, color=ACCENT_BLUE, bold=True)

    # Key insight
    box = add_shape_rect(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(1.5),
                         RGBColor(0x1E, 0x2D, 0x50), ACCENT_ORANGE)
    add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(1),
                 "⚡ L'ajout d'un vecteur est rapide : il suffit de comparer avec les 'nlist' "
                 "centroïdes (au lieu de tous les vecteurs), puis d'ajouter l'ID à la posting list.",
                 font_size=15, color=TEXT_YELLOW)


def create_retrieval_overview_slide(prs):
    """Slide 7 : Vue d'ensemble de la phase de recherche."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "IVF - Phase de Recherche (Retrieval)",
                 font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_PURPLE)

    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.6),
                 "Comment l'IVF trouve rapidement les voisins les plus proches",
                 font_size=18, color=TEXT_LIGHT)

    # Les 3 étapes de la recherche
    steps = [
        ("1", "Trouver les centroïdes\nles plus proches",
         "Pour un vecteur requête, calculer\nla distance aux centroïdes et\nsélectionner les 'nprobe'\nplus proches.",
         ACCENT_BLUE),
        ("2", "Explorer les clusters\nsélectionnés",
         "Chercher les voisins les plus\nproches uniquement dans\nles 'nprobe' clusters\nsélectionnés.",
         ACCENT_GREEN),
        ("3", "Retourner les top-k\nrésultats",
         "Fusionner les résultats de tous\nles clusters explorés et retourner\nles k vecteurs les plus proches.",
         ACCENT_ORANGE),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.8 + i * 4.1)
        y = Inches(2.6)

        # Numéro
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Titre
        add_text_box(slide, x, y + Inches(1), Inches(3.8), Inches(0.8),
                     title, font_size=17, color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, x, y + Inches(1.9), Inches(3.8), Inches(1.5),
                     desc, font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


def create_nprobe_slide(prs):
    """Slide 8 : Le paramètre nprobe en détail."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Le Paramètre nprobe", font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_BLUE)

    # Explication
    add_bullet_points(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(1.5), [
        "nprobe définit le nombre de clusters à explorer lors d'une requête",
        "C'est le paramètre clé pour contrôler le compromis vitesse/précision",
    ], font_size=16, color=TEXT_LIGHT)

    # Comparaison nprobe faible vs élevé
    # Low nprobe
    box1 = add_shape_rect(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(3.5),
                          RGBColor(0x1E, 0x2D, 0x50), ACCENT_ORANGE)
    add_text_box(slide, Inches(1.2), Inches(3.4), Inches(4.8), Inches(0.5),
                 "nprobe faible (ex: 1-2)", font_size=20, color=ACCENT_ORANGE, bold=True)
    add_bullet_points(slide, Inches(1.2), Inches(4), Inches(4.8), Inches(2.5), [
        "✅ Recherche très rapide",
        "✅ Faible latence",
        "❌ Recall potentiellement bas",
        "❌ Peut manquer des résultats pertinents",
    ], font_size=14, color=TEXT_LIGHT)

    # High nprobe
    box2 = add_shape_rect(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.5),
                          RGBColor(0x1E, 0x2D, 0x50), ACCENT_GREEN)
    add_text_box(slide, Inches(7.2), Inches(3.4), Inches(4.8), Inches(0.5),
                 "nprobe élevé (ex: nlist)", font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_points(slide, Inches(7.2), Inches(4), Inches(4.8), Inches(2.5), [
        "✅ Recall élevé (proche de 100%)",
        "✅ Résultats très précis",
        "❌ Recherche plus lente",
        "❌ Latence plus élevée",
    ], font_size=14, color=TEXT_LIGHT)


def create_parameters_slide(prs):
    """Slide 9 : Paramètres clés de l'IVF."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Paramètres Clés de l'IVF", font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_BLUE)

    # nlist parameter
    box1 = add_shape_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                          RGBColor(0x1E, 0x2D, 0x50), ACCENT_BLUE)
    add_text_box(slide, Inches(1.2), Inches(2), Inches(4.8), Inches(0.6),
                 "nlist", font_size=28, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(0.5),
                 "Nombre de clusters", font_size=18, color=TEXT_WHITE)
    add_bullet_points(slide, Inches(1.2), Inches(3.2), Inches(4.8), Inches(2.5), [
        "Détermine combien de clusters sont créés lors de l'indexation",
        "Plus de clusters = posting lists plus petites",
        "Valeur typique : nlist = sqrt(n)",
        "Influence la granularité de la partition",
        "Doit être ajusté selon la taille du dataset",
    ], font_size=13, color=TEXT_LIGHT)

    # nprobe parameter
    box2 = add_shape_rect(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5),
                          RGBColor(0x1E, 0x2D, 0x50), ACCENT_GREEN)
    add_text_box(slide, Inches(7.2), Inches(2), Inches(4.8), Inches(0.6),
                 "nprobe", font_size=28, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(7.2), Inches(2.6), Inches(4.8), Inches(0.5),
                 "Nombre de clusters à explorer", font_size=18, color=TEXT_WHITE)
    add_bullet_points(slide, Inches(7.2), Inches(3.2), Inches(4.8), Inches(2.5), [
        "Nombre de clusters explorés lors d'une requête",
        "nprobe élevé = meilleur recall, latence plus haute",
        "nprobe faible = recherche rapide, recall moindre",
        "Si nprobe = nlist → recherche exhaustive",
        "Contrôle le compromis vitesse/qualité",
    ], font_size=13, color=TEXT_LIGHT)

    # K parameter
    box3 = add_shape_rect(slide, Inches(3.5), Inches(6.5), Inches(6), Inches(0.7),
                          RGBColor(0x1E, 0x2D, 0x50), ACCENT_ORANGE)
    add_text_box(slide, Inches(3.7), Inches(6.55), Inches(5.6), Inches(0.5),
                 "k : nombre de résultats les plus proches à retourner (top-k)",
                 font_size=14, color=ACCENT_ORANGE)


def create_summary_slide(prs):
    """Slide 10 : Récapitulatif."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Récapitulatif : IVF en un coup d'œil",
                 font_size=36, color=TEXT_WHITE, bold=True)

    add_shape_rect(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.03), ACCENT_BLUE)

    # Indexation
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
                 "📥  Indexation", font_size=22, color=ACCENT_BLUE, bold=True)
    add_bullet_points(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2), [
        "Clustering K-means → 'nlist' clusters",
        "Chaque cluster → centroïde + posting list d'IDs",
        "Nouveau vecteur → ajouté au cluster le plus proche",
    ], font_size=14, color=TEXT_LIGHT)

    # Recherche
    add_text_box(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.5),
                 "🔍  Recherche", font_size=22, color=ACCENT_GREEN, bold=True)
    add_bullet_points(slide, Inches(6.8), Inches(2.3), Inches(5.5), Inches(2), [
        "Trouver les 'nprobe' centroïdes les plus proches",
        "Explorer uniquement ces 'nprobe' clusters",
        "Retourner les top-k résultats combinés",
    ], font_size=14, color=TEXT_LIGHT)

    # Avantages et limites
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.5),
                 "✅  Avantages", font_size=20, color=ACCENT_GREEN, bold=True)
    add_bullet_points(slide, Inches(0.8), Inches(5), Inches(5.5), Inches(2), [
        "Recherche beaucoup plus rapide que brute force",
        "Compromis vitesse/précision contrôlable",
        "Simple à comprendre et implémenter",
        "Fonctionne bien avec FAISS et autres bibliothèques",
    ], font_size=13, color=TEXT_LIGHT)

    add_text_box(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(0.5),
                 "⚠️  Limites", font_size=20, color=ACCENT_ORANGE, bold=True)
    add_bullet_points(slide, Inches(6.8), Inches(5), Inches(5.5), Inches(2), [
        "Recherche approximative (pas exacte)",
        "La qualité dépend du clustering initial",
        "Nécessite un bon réglage de nlist et nprobe",
        "Le clustering K-means peut être coûteux sur de gros datasets",
    ], font_size=13, color=TEXT_LIGHT)


def main():
    """Crée la présentation PowerPoint complète sur l'algorithme IVF."""
    print("=" * 60)
    print("Création de la présentation IVF")
    print("=" * 60 + "\n")

    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Création des slides...")
    create_title_slide(prs)
    print("  ✓ Slide 1 : Titre")

    create_intro_slide(prs)
    print("  ✓ Slide 2 : Introduction - Pourquoi IVF ?")

    create_indexing_overview_slide(prs)
    print("  ✓ Slide 3 : Vue d'ensemble de l'indexation")

    create_clustering_slide(prs)
    print("  ✓ Slide 4 : Clustering K-means")

    create_posting_list_slide(prs)
    print("  ✓ Slide 5 : Centroïdes & Posting Lists")

    create_adding_vector_slide(prs)
    print("  ✓ Slide 6 : Ajout d'un nouveau vecteur")

    create_retrieval_overview_slide(prs)
    print("  ✓ Slide 7 : Phase de recherche")

    create_nprobe_slide(prs)
    print("  ✓ Slide 8 : Le paramètre nprobe")

    create_parameters_slide(prs)
    print("  ✓ Slide 9 : Paramètres clés")

    create_summary_slide(prs)
    print("  ✓ Slide 10 : Récapitulatif")

    prs.save(PPTX_PATH)
    print(f"\n✅ Présentation sauvegardée : {PPTX_PATH}")
    print(f"   ({len(prs.slides)} slides)")
    print("\n💡 Pour importer dans Canva :")
    print("   1. Aller sur https://www.canva.com/")
    print("   2. Cliquer sur 'Créer un design' → 'Importer'")
    print("   3. Sélectionner le fichier presentation_ivf.pptx")
    print("   4. Canva convertira automatiquement le fichier en design éditable")


if __name__ == "__main__":
    main()
